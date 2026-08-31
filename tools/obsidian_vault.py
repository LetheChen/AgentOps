"""obsidian_vault 工具：Obsidian vault 统一访问层（Layer A）。

背景：两个个人助理 agent（content_curator / proposal_planner）
都要读写 E:\\Document Obsidian vault。若每个 agent 各写各的会路径硬编码、
frontmatter 不一致、非 md 文件无法处理。本工具作为 Layer A 统一访问入口，
封装 path 白名单 + frontmatter 自动注入 + 非 md 文件文本抽取。

设计：
- 7 个 action：list_files / read_file / write_file / search_by_keyword / search_by_tag / scan_incremental / validate_path
- 读操作可读全 vault，写操作只能在 Articles\\Reports\\Notes\\Weekly\\Images\\ 下（保护用户原始笔记）
- write_file 自动注入 frontmatter（ingested_at / ingested_by / type / tags）
- read_file 对非 md 文件（pdf/pptx/docx/xlsx/html）统一文本抽取，lazy import 避免强依赖
- 命名规则校验（如 Weekly\\ 下必须匹配 工作周报_YYYY-MM-DD_MM-DD.md）

参考配置：config/tools/obsidian_vault.yaml
设计文档：docs/knowledge-base/DESIGN_knowledge_base_module.md §3
"""
from __future__ import annotations

import logging
import os
import re
from datetime import datetime, timezone
from pathlib import Path
from typing import Any

logger = logging.getLogger(__name__)

# Vault 根目录：解析优先级 环境变量 OBSIDIAN_VAULT_ROOT > config/tools/obsidian_vault.yaml
# 的 handler.config.vault_root > 硬编码默认（向后兼容本地开发）。
# 生产/容器环境通常不挂载个人 vault，建议设 OBSIDIAN_VAULT_ROOT="" 显式禁用。
_VAULT_ROOT_DEFAULT = Path(r"E:\Document")


def _load_vault_root() -> Path:
    """解析 vault_root：环境变量 > yaml > 硬编码默认。

    环境变量 OBSIDIAN_VAULT_ROOT 存在时（包括空字符串）直接采用：
    - 非空 → 作为 vault 根目录
    - 空串 → 显式禁用个人 vault（生产/容器部署），返回不存在的路径占位，
      所有 vault 读写会因目录缺失而报错/返回空，不会误访问他人机器目录。
    """
    env_val = os.environ.get("OBSIDIAN_VAULT_ROOT")
    if env_val is not None:
        if env_val == "":
            logger.info("obsidian_vault 经 OBSIDIAN_VAULT_ROOT='' 显式禁用（个人 vault 不随生产部署）")
            return Path(os.devnull) / "obsidian_vault_disabled"
        logger.info("obsidian_vault 从环境变量 OBSIDIAN_VAULT_ROOT 加载 vault_root=%s", env_val)
        return Path(env_val)
    import yaml
    config_path = Path(__file__).resolve().parent.parent / "config" / "tools" / "obsidian_vault.yaml"
    try:
        if config_path.exists():
            with open(config_path, encoding="utf-8") as f:
                cfg = yaml.safe_load(f) or {}
            handler_cfg = cfg.get("handler", {}).get("config", {})
            vault_str = handler_cfg.get("vault_root", "")
            if vault_str:
                logger.info("obsidian_vault 从 yaml 加载 vault_root=%s", vault_str)
                return Path(vault_str)
    except Exception as e:
        logger.warning("obsidian_vault 加载 yaml 失败，回退硬编码: %s", e)
    logger.info("obsidian_vault 使用硬编码 vault_root=%s", _VAULT_ROOT_DEFAULT)
    return _VAULT_ROOT_DEFAULT


VAULT_ROOT = _load_vault_root()

# 写操作白名单（只能在这些目录下创建/覆盖文件）
WRITE_WHITELIST = [
    "Articles\\",
    "Reports\\",
    "Notes\\",
    "Weekly\\",
    "Images\\",
]

# 命名规则（按子目录配置）
NAMING_RULES = {
    "Weekly\\": {
        "pattern": re.compile(r"^工作周报_\d{4}-\d{2}-\d{2}_\d{2}-\d{2}\.md$"),
        "description": "周报命名规则：工作周报_YYYY-MM-DD_MM-DD.md",
    },
}

# 归档目录 → type 字段映射（用于 frontmatter 自动注入）
PATH_TYPE_MAP = {
    "Articles\\": "article",
    "Reports\\": "report",
    "Notes\\": "note",
    "Weekly\\": "content_curator",
    "Images\\": "image",
}

# 支持文本抽取的扩展名 → 抽取器名称
EXTRACTORS = {
    "md": "raw",
    "txt": "raw",
    "xml": "raw",
    "yaml": "raw",
    "yml": "raw",
    "json": "raw",
    "pdf": "pdfplumber",
    "pptx": "python-pptx",
    "docx": "python-docx",
    "xlsx": "openpyxl",
    "xls": "openpyxl",
    "html": "beautifulsoup4",
    "htm": "beautifulsoup4",
}


async def obsidian_vault(args: dict[str, Any]) -> dict[str, Any]:
    """Obsidian vault 统一访问工具。

    Args:
        args: dict，必须包含 action 字段，其他字段随 action 变化
            - action (str, required): list_files | read_file | write_file | search_by_keyword | search_by_tag | scan_incremental | validate_path
            - path (str): vault 内相对路径
            - content (str): write_file 时的内容
            - frontmatter (dict): write_file 时的 frontmatter
            - ext_filter (list[str]): list_files 时的扩展名过滤
            - tag_filter (list[str]): list_files/search_by_tag 时的 tag 过滤
            - query (str): search_by_keyword 时的搜索词
            - since (str): scan_incremental 时的起始时间（ISO 8601）
            - max_results (int): 搜索结果上限，默认 100
            - agent_id (str): 调用方 agent id（用于 frontmatter ingested_by 字段）

    Returns:
        dict，至少包含 content（人类可读摘要）+ action 特定字段。
        失败时包含 error 字段。
    """
    action = (args.get("action") or "").strip()
    if not action:
        return {"content": "调用失败：缺少 action 参数", "error": "missing_action"}

    try:
        if action == "list_files":
            return _action_list_files(args)
        elif action == "read_file":
            return _action_read_file(args)
        elif action == "write_file":
            return await _action_write_file(args)
        elif action == "search_by_keyword":
            return _action_search_by_keyword(args)
        elif action == "search_by_tag":
            return _action_search_by_tag(args)
        elif action == "scan_incremental":
            return _action_scan_incremental(args)
        elif action == "validate_path":
            return _action_validate_path(args)
        else:
            return {
                "content": f"调用失败：未知 action '{action}'",
                "error": "unknown_action",
                "available_actions": [
                    "list_files", "read_file", "write_file",
                    "search_by_keyword", "search_by_tag",
                    "scan_incremental", "validate_path",
                ],
            }
    except Exception as e:
        logger.exception("obsidian_vault action=%s failed: %s", action, e)
        return {"content": f"调用失败：{e}", "error": "action_failed", "action": action}


# ==================== action 实现 ====================


def _action_list_files(args: dict[str, Any]) -> dict[str, Any]:
    """列举 vault 内指定路径下的文件。"""
    rel_path = (args.get("path") or "").strip()
    ext_filter = args.get("ext_filter") or []
    tag_filter = args.get("tag_filter") or []
    max_results = int(args.get("max_results") or 100)

    target_dir = _resolve_vault_path(rel_path)
    if not target_dir.exists():
        return {"content": f"路径不存在：{rel_path}", "error": "path_not_found"}
    if not target_dir.is_dir():
        return {"content": f"不是目录：{rel_path}", "error": "not_a_directory"}

    # 收集文件
    files = []
    for f in target_dir.rglob("*"):
        if not f.is_file():
            continue
        if ext_filter and f.suffix.lstrip(".").lower() not in [e.lower() for e in ext_filter]:
            continue
        # 读 frontmatter 做 tag 过滤
        fm = _parse_frontmatter(f) if f.suffix.lower() == ".md" else {}
        if tag_filter and not _match_tags(fm.get("tags", []), tag_filter):
            continue
        files.append({
            "path": str(f.relative_to(VAULT_ROOT)),
            "size": f.stat().st_size,
            "mtime": datetime.fromtimestamp(f.stat().st_mtime, tz=timezone.utc).isoformat(),
            "tags": fm.get("tags", []),
        })
        if len(files) >= max_results:
            break

    return {
        "content": f"在 {rel_path} 下找到 {len(files)} 个文件",
        "files": files,
        "total": len(files),
        "truncated": len(files) >= max_results,
    }


def _action_read_file(args: dict[str, Any]) -> dict[str, Any]:
    """读取文件，非 md 文件自动文本抽取。"""
    rel_path = (args.get("path") or "").strip()
    if not rel_path:
        return {"content": "调用失败：缺少 path 参数", "error": "missing_path"}

    file_path = _resolve_vault_path(rel_path)
    if not file_path.exists():
        return {"content": f"文件不存在：{rel_path}", "error": "file_not_found"}
    if not file_path.is_file():
        return {"content": f"不是文件：{rel_path}", "error": "not_a_file"}

    ext = file_path.suffix.lstrip(".").lower()
    extractor = EXTRACTORS.get(ext, "raw")

    try:
        content_text, extracted_by = _extract_text(file_path, extractor)
    except ImportError as e:
        return {
            "content": f"文本抽取失败：缺少依赖 {e.name}（扩展名 .{ext}）",
            "error": "missing_dependency",
            "missing_module": e.name,
            "ext": ext,
        }
    except Exception as e:
        logger.exception("read_file 抽取失败 path=%s: %s", rel_path, e)
        return {"content": f"文本抽取失败：{e}", "error": "extract_failed"}

    return {
        "content": content_text[:500] + ("..." if len(content_text) > 500 else ""),
        "full_content": content_text,
        "path": rel_path,
        "format": "markdown",
        "extracted_by": extracted_by,
        "size": file_path.stat().st_size,
    }


async def _action_write_file(args: dict[str, Any]) -> dict[str, Any]:
    """写入文件，path 白名单 + frontmatter 自动注入 + 命名规则校验。"""
    rel_path = (args.get("path") or "").strip()
    content = args.get("content") or ""
    frontmatter = args.get("frontmatter") or {}
    agent_id = args.get("agent_id") or "unknown_agent"

    if not rel_path:
        return {"content": "调用失败：缺少 path 参数", "error": "missing_path"}

    # path 白名单校验
    validation = _validate_write_path(rel_path)
    if not validation["write_allowed"]:
        return {
            "content": f"写入被拒绝：{validation['reason']}",
            "error": "write_denied",
            "reason": validation["reason"],
        }

    file_path = _resolve_vault_path(rel_path)

    # 命名规则校验
    naming_error = _check_naming_rule(rel_path)
    if naming_error:
        return {
            "content": f"命名规则不匹配：{naming_error}",
            "error": "naming_violation",
            "rule": naming_error,
        }

    # 确保父目录存在
    file_path.parent.mkdir(parents=True, exist_ok=True)

    # 注入自动 frontmatter 字段（auto 字段优先级低，用户传入的覆盖 auto 的）
    auto_fm = _build_auto_frontmatter(rel_path, agent_id)
    final_fm = {**auto_fm, **frontmatter}

    existed = file_path.exists()
    action_type = "overwritten" if existed else "created"

    # 组装最终内容（frontmatter + 正文）
    final_content = _serialize_with_frontmatter(final_fm, content)

    file_path.write_text(final_content, encoding="utf-8")
    logger.info(
        "obsidian_vault write_file path=%s action=%s size=%d agent=%s",
        rel_path, action_type, len(final_content), agent_id,
    )

    return {
        "content": f"文件已{action_type}：{rel_path}（{len(final_content)} 字节）",
        "path": rel_path,
        "absolute_path": str(file_path),
        "size": len(final_content),
        "frontmatter_injected": list(auto_fm.keys()),
        "existed": existed,
        "action": action_type,
    }


def _action_search_by_keyword(args: dict[str, Any]) -> dict[str, Any]:
    """基于 ripgrep 的全文搜索（仅 md 文件）。"""
    query = (args.get("query") or "").strip()
    max_results = int(args.get("max_results") or 100)

    if not query:
        return {"content": "调用失败：缺少 query 参数", "error": "missing_query"}

    # 用 Grep 工具的逻辑（这里直接用 re 扫描，避免依赖外部 rg）
    matches = []
    pattern = re.compile(re.escape(query), re.IGNORECASE)

    for md_file in VAULT_ROOT.rglob("*.md"):
        try:
            text = md_file.read_text(encoding="utf-8", errors="ignore")
        except OSError:
            continue
        rel = str(md_file.relative_to(VAULT_ROOT))
        for line_no, line in enumerate(text.split("\n"), start=1):
            if pattern.search(line):
                # 上下文（前后各 20 字符）
                idx = line.lower().find(query.lower())
                start = max(0, idx - 20)
                end = min(len(line), idx + len(query) + 20)
                context = ("..." if start > 0 else "") + line[start:end] + ("..." if end < len(line) else "")
                matches.append({
                    "path": rel,
                    "line": line_no,
                    "context": context,
                })
                if len(matches) >= max_results:
                    break
        if len(matches) >= max_results:
            break

    return {
        "content": f"找到 {len(matches)} 处匹配 '{query}'",
        "matches": matches,
        "total": len(matches),
        "truncated": len(matches) >= max_results,
    }


def _action_search_by_tag(args: dict[str, Any]) -> dict[str, Any]:
    """按 frontmatter tag 搜索。"""
    tags = args.get("tag_filter") or []
    if isinstance(tags, str):
        tags = [tags]
    max_results = int(args.get("max_results") or 100)

    if not tags:
        return {"content": "调用失败：缺少 tag_filter 参数", "error": "missing_tag_filter"}

    matches = []
    for md_file in VAULT_ROOT.rglob("*.md"):
        fm = _parse_frontmatter(md_file)
        file_tags = fm.get("tags", [])
        if isinstance(file_tags, str):
            file_tags = [file_tags]
        if _match_tags(file_tags, tags):
            matches.append({
                "path": str(md_file.relative_to(VAULT_ROOT)),
                "size": md_file.stat().st_size,
                "mtime": datetime.fromtimestamp(md_file.stat().st_mtime, tz=timezone.utc).isoformat(),
                "tags": file_tags,
            })
            if len(matches) >= max_results:
                break

    return {
        "content": f"找到 {len(matches)} 个匹配 tag {tags} 的文件",
        "matches": matches,
        "total": len(matches),
        "truncated": len(matches) >= max_results,
    }


def _action_scan_incremental(args: dict[str, Any]) -> dict[str, Any]:
    """增量扫描 vault 内 mtime > since 的文件。"""
    since_str = (args.get("since") or "").strip()
    if not since_str:
        return {"content": "调用失败：缺少 since 参数", "error": "missing_since"}

    try:
        since_dt = datetime.fromisoformat(since_str)
        if since_dt.tzinfo is None:
            since_dt = since_dt.replace(tzinfo=timezone.utc)
    except ValueError as e:
        return {"content": f"since 解析失败：{e}", "error": "invalid_since"}

    since_ts = since_dt.timestamp()
    new_files = []
    modified_files = []

    for f in VAULT_ROOT.rglob("*"):
        if not f.is_file():
            continue
        # 跳过 .obsidian/ 配置目录
        if ".obsidian" in f.parts:
            continue
        mtime = f.stat().st_mtime
        if mtime <= since_ts:
            continue
        rel = str(f.relative_to(VAULT_ROOT))
        # 判断 new vs modified：基于 frontmatter.ingested_at 是否存在
        is_new = True
        if f.suffix.lower() == ".md":
            fm = _parse_frontmatter(f)
            if fm.get("ingested_at"):
                is_new = False
        entry = {
            "path": rel,
            "size": f.stat().st_size,
            "mtime": datetime.fromtimestamp(mtime, tz=timezone.utc).isoformat(),
        }
        if is_new:
            new_files.append(entry)
        else:
            modified_files.append(entry)

    return {
        "content": f"扫描完成：{len(new_files)} 个新文件 + {len(modified_files)} 个修改文件",
        "new_files": new_files,
        "modified_files": modified_files,
        "scanned_at": datetime.now(timezone.utc).isoformat(),
    }


def _action_validate_path(args: dict[str, Any]) -> dict[str, Any]:
    """路径白名单校验，返回是否允许读/写。"""
    rel_path = (args.get("path") or "").strip()
    if not rel_path:
        return {"content": "调用失败：缺少 path 参数", "error": "missing_path"}

    read_allowed = True
    write_validation = _validate_write_path(rel_path)
    reason = write_validation["reason"] if not write_validation["write_allowed"] else "within write whitelist"

    return {
        "content": f"path={rel_path} read={read_allowed} write={write_validation['write_allowed']}（{reason}）",
        "path": rel_path,
        "read_allowed": read_allowed,
        "write_allowed": write_validation["write_allowed"],
        "reason": reason,
    }


# ==================== 辅助函数 ====================


def _resolve_vault_path(rel_path: str) -> Path:
    """把 vault 内相对路径解析为绝对路径，防止路径遍历。"""
    # 标准化分隔符
    normalized = rel_path.replace("/", "\\").lstrip(".\\")
    abs_path = (VAULT_ROOT / normalized).resolve()
    # 校验仍在 vault 内
    try:
        abs_path.relative_to(VAULT_ROOT.resolve())
    except ValueError:
        raise ValueError(f"路径越界：{rel_path} 不在 vault 内")
    return abs_path


def _validate_write_path(rel_path: str) -> dict[str, bool, str]:
    """校验写路径是否在白名单内。"""
    normalized = rel_path.replace("/", "\\")
    for whitelist_dir in WRITE_WHITELIST:
        if normalized.lower().startswith(whitelist_dir.lower()):
            return {"write_allowed": True, "reason": "within write whitelist"}
    return {
        "write_allowed": False,
        "reason": f"只能写入 Articles/Reports/Notes/Weekly/Images/ 下，当前路径：{rel_path}",
    }


def _check_naming_rule(rel_path: str) -> str | None:
    """检查路径是否符合命名规则，返回错误描述或 None。"""
    normalized = rel_path.replace("/", "\\")
    for whitelist_dir, rule in NAMING_RULES.items():
        if normalized.lower().startswith(whitelist_dir.lower()):
            filename = Path(normalized).name
            if not rule["pattern"].match(filename):
                return rule["description"]
    return None


def _build_auto_frontmatter(rel_path: str, agent_id: str) -> dict[str, Any]:
    """构建自动注入的 frontmatter 字段。"""
    normalized = rel_path.replace("/", "\\")
    file_type = "note"  # 默认
    for path_prefix, type_val in PATH_TYPE_MAP.items():
        if normalized.lower().startswith(path_prefix.lower()):
            file_type = type_val
            break

    return {
        "ingested_at": datetime.now(timezone.utc).isoformat(),
        "ingested_by": agent_id,
        "type": file_type,
        "tags": ["agent-generated"],
    }


def _parse_frontmatter(file_path: Path) -> dict[str, Any]:
    """解析 md 文件的 YAML frontmatter。"""
    if not file_path.exists() or file_path.suffix.lower() != ".md":
        return {}
    try:
        text = file_path.read_text(encoding="utf-8", errors="ignore")
    except OSError:
        return {}
    if not text.startswith("---"):
        return {}
    # 找第二个 --- 作为 frontmatter 结束
    end_idx = text.find("\n---", 3)
    if end_idx == -1:
        return {}
    fm_text = text[3:end_idx].strip()
    try:
        import yaml
        return yaml.safe_load(fm_text) or {}
    except (yaml.YAMLError, ImportError):
        return {}


def _match_tags(file_tags: list | str, filter_tags: list[str]) -> bool:
    """判断文件 tags 是否匹配过滤条件（任一匹配即可）。"""
    if isinstance(file_tags, str):
        file_tags = [file_tags]
    file_tags_lower = [str(t).lower() for t in file_tags]
    filter_lower = [t.lower() for t in filter_tags]
    return any(ft in file_tags_lower for ft in filter_lower)


def _serialize_with_frontmatter(frontmatter: dict[str, Any], content: str) -> str:
    """组装 frontmatter + 正文。"""
    try:
        import yaml
        # tags 特殊处理：确保是 list
        if "tags" in frontmatter and isinstance(frontmatter["tags"], str):
            frontmatter["tags"] = [frontmatter["tags"]]
        fm_yaml = yaml.safe_dump(frontmatter, allow_unicode=True, default_flow_style=False, sort_keys=False).strip()
        return f"---\n{fm_yaml}\n---\n\n{content}"
    except (ImportError, yaml.YAMLError):
        # fallback：无 yaml 也能写
        fm_lines = []
        for k, v in frontmatter.items():
            if isinstance(v, list):
                fm_lines.append(f"{k}:")
                for item in v:
                    fm_lines.append(f"  - {item}")
            else:
                fm_lines.append(f"{k}: {v}")
        return f"---\n{chr(10).join(fm_lines)}\n---\n\n{content}"


def _extract_text(file_path: Path, extractor: str) -> tuple[str, str]:
    """根据抽取器名称抽取文本，返回 (text, extracted_by)。lazy import 避免强依赖。"""
    if extractor == "raw":
        return file_path.read_text(encoding="utf-8", errors="ignore"), "raw"

    if extractor == "pdfplumber":
        try:
            import pdfplumber
        except ImportError:
            raise ImportError("pdfplumber")
        pages = []
        with pdfplumber.open(file_path) as pdf:
            for i, page in enumerate(pdf.pages, start=1):
                text = page.extract_text() or ""
                pages.append(f"## 第 {i} 页\n\n{text}")
        return "\n\n".join(pages), "pdfplumber"

    if extractor == "python-pptx":
        try:
            from pptx import Presentation
        except ImportError:
            raise ImportError("python-pptx")
        slides = []
        prs = Presentation(file_path)
        for i, slide in enumerate(prs.slides, start=1):
            texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        line = "".join(run.text for run in para.runs)
                        if line.strip():
                            texts.append(line)
            slides.append(f"## Slide {i}\n\n" + "\n".join(texts))
        return "\n\n".join(slides), "python-pptx"

    if extractor == "python-docx":
        try:
            import docx
        except ImportError:
            raise ImportError("python-docx")
        doc = docx.Document(file_path)
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        return "\n\n".join(paragraphs), "python-docx"

    if extractor == "openpyxl":
        try:
            from openpyxl import load_workbook
        except ImportError:
            raise ImportError("openpyxl")
        wb = load_workbook(file_path, read_only=True, data_only=True)
        sheets = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            rows = []
            for row in ws.iter_rows(values_only=True):
                row_text = " | ".join(str(c) if c is not None else "" for c in row)
                if row_text.strip(" |"):
                    rows.append(f"| {row_text} |")
            sheets.append(f"## Sheet: {sheet_name}\n\n" + "\n".join(rows))
        wb.close()
        return "\n\n".join(sheets), "openpyxl"

    if extractor == "beautifulsoup4":
        try:
            from bs4 import BeautifulSoup
        except ImportError:
            raise ImportError("beautifulsoup4")
        html = file_path.read_text(encoding="utf-8", errors="ignore")
        soup = BeautifulSoup(html, "html.parser")
        # 去掉 script/style
        for tag in soup(["script", "style"]):
            tag.decompose()
        return soup.get_text(separator="\n", strip=True), "beautifulsoup4"

    # 兜底：按 raw 处理
    return file_path.read_text(encoding="utf-8", errors="ignore"), "raw"
